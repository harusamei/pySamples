# python-pptx 用法可参考 https://zhuanlan.zhihu.com/p/265464562
# 基本的layout有9种，https://python-pptx.readthedocs.io/en/latest/user/slides.html?highlight=layout#slide-layout-basics
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
from copy import deepcopy
import glob
import sys

def getContent(shape):
    content = ""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            content += getContent(sub_shape)
    elif shape.has_text_frame:
        content += shape.text_frame.text
        content += "\n"
    return content

def get_slide_content(slide):
    content = ""
    for shape in slide.shapes:
        content += getContent(shape)
    return content

def edit_slide(slide):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            shape.line.width = Inches(0.2)
            shape.line.color.rgb = RGBColor(255, 0, 0)
    return slide

def add_group_shape(shape,slide):
    
    for sub_shape in shape.shapes:
        #print(f"sub_shape {sub_shape.shape_type} name {sub_shape.name}")
    
        if sub_shape.left is not None:
            sub_shape.left += shape.left
            sub_shape.top += shape.top

        if sub_shape.shape_type != MSO_SHAPE_TYPE.GROUP:
            add_nogroup_shape(sub_shape, slide)
        else:
            add_group_shape(sub_shape, slide)
    return slide

def add_nogroup_shape(shape, slide):
    # print(f"type {shape.shape_type} name {shape.name}")
    if shape.shape_type == MSO_SHAPE_TYPE.LINE:
        new_shape = slide.shapes.add_shape( MSO_SHAPE_TYPE.LINE, shape.left, shape.top, shape.width, shape.height )
        new_shape.line.color.rgb = RGBColor(255, 0, 0) #shape.line.color.rgb  # 设置线条红色
        new_shape.line.width = shape.line.width          # 设置线条宽度
    elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        new_shape= slide.shapes.add_shape(MSO_SHAPE_TYPE.AUTO_SHAPE,shape.left, shape.top, shape.width, shape.height)
        new_shape.text_frame.text = shape.text
    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        # print(f"shape.text {shape.text} [{shape.left/914400} {shape.top/914400} {shape.width} {shape.height} ]")
        if shape.width < Inches(0.5): # 针对sub-shape的absolute width, height不知道怎么计算
            shape.width = Inches(1) 
            shape.height = Inches(1)
        new_shape = slide.shapes.add_textbox(shape.left, shape.top ,shape.width, shape.height)
        new_shape.text_frame.text = shape.text
    else:
        new_element = deepcopy(shape.element)
        slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')
    return slide

def add_slide(slide, prs):

    new_slide = prs.slides.add_slide(prs.slide_layouts[27])  # 6 blank layout, 1 title content
    print(f" num of shape {len(slide.shapes)}") 
    for i, shape in enumerate(slide.shapes):
        print(f"shape {i} type {shape.shape_type} name {shape.name}")
        print(f"[ {shape.left} {shape.top} {shape.width} {shape.height} ]")
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            add_group_shape(shape, new_slide)
        else:
            add_nogroup_shape(shape, new_slide)

    return new_slide

if __name__ == '__main__':
    
    file_names = glob.glob('data/*.pptx')
    # 创建一个新的PowerPoint文件
    mergedPrs = Presentation('data/z.pptx')

    # 遍历每个文件并将其内容复制到新的文件中
    for file_name in file_names:
        if file_name in ['data\merged.pptx','data\z.pptx']:
            continue
        onePPT = Presentation(file_name)
        print(f"Processing {file_name}, {len(onePPT.slides)}")
        new_slide = mergedPrs.slides.add_slide(mergedPrs.slide_layouts[0])
        new_slide.shapes.title.text = file_name
        for slide in onePPT.slides:
            new_slide = add_slide(slide, mergedPrs)
        
            
    # 保存合并后的PowerPoint文件
    mergedPrs.save('data/merged.pptx')
