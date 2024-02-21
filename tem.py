from pptx import Presentation
import glob
from copy import deepcopy

def copy_slide(s, t):
    # empty the target slide
    for shape in t.shapes:
        sp = shape.element
        t.shapes._spTree.remove(sp)
        
    # Copy the content from the source slide to the target slide
    for shape in s.shapes:
        new_shape = deepcopy(shape.element)
        t.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

    return t

# List all PPT files in the directory
ppt_files = glob.glob("data\*.pptx")

# Create a new presentation object
allPPT = Presentation()

# Iterate over each PPT file
for ppt_file in ppt_files:
    # Open the PPT file
    onePPT = Presentation(ppt_file)
    print(len(onePPT.slides))

    # Iterate over each slide in the PPT file
    for slide in onePPT.slides:
        # 增加一页空白页
        tSlide= allPPT.slides.add_slide(allPPT.slide_layouts[5])
        copy_slide(slide, tSlide)

# Save the merged presentation to a new file
allPPT.save("data\merged.pptx")
